"""
Premium PowerPoint export for executive intelligence reports.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import APP_NAME
from services.export_chart_blocks import get_export_chart_images
from services.report_document_parser import (
    dashboard_metrics,
    first_ai_insight,
    is_available,
    parse_intelligence_report,
    top_risks,
)

BLUE = RGBColor(29, 78, 216)
SLATE = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)


@dataclass
class PresentationExportMetadata:
    project_name: str
    report_name: str
    source_documents: list[str] | None = None
    # "Branded reports with your logo" (Professional+): see
    # PremiumExportMetadata.custom_logo_bytes for the same field on the PDF
    # builder.
    custom_logo_bytes: bytes | None = None
    # "Prepared by" attribution line — see PremiumExportMetadata.prepared_by
    # for the same field on the PDF builder. Falls back to APP_NAME when blank.
    prepared_by: str | None = None


class PremiumPresentationBuilder:
    def __init__(self, metadata: PresentationExportMetadata) -> None:
        self.metadata = metadata
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

    def _title_slide(self) -> None:
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        if self.metadata.custom_logo_bytes:
            slide.shapes.add_picture(
                BytesIO(self.metadata.custom_logo_bytes),
                Inches(5.92),
                Inches(0.4),
                height=Inches(0.7),
            )

        title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1))
        frame = title.text_frame
        frame.text = APP_NAME
        frame.paragraphs[0].font.size = Pt(34)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = SLATE
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.8))
        sub_frame = subtitle.text_frame
        sub_frame.text = "Executive Intelligence Presentation"
        sub_frame.paragraphs[0].font.size = Pt(20)
        sub_frame.paragraphs[0].font.color.rgb = BLUE
        sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        meta = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(2))
        meta_frame = meta.text_frame
        meta_frame.text = (
            f"Prepared by: {self.metadata.prepared_by or APP_NAME}\n"
            f"Project: {self.metadata.project_name}\n"
            f"Report: {self.metadata.report_name}\n"
            f"Generated: {datetime.now(timezone.utc).strftime('%d %B %Y')}"
        )
        for paragraph in meta_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = MUTED
            paragraph.alignment = PP_ALIGN.CENTER

        disclaimer = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5))
        disclaimer_frame = disclaimer.text_frame
        disclaimer_frame.text = (
            "DataDumpAI can make mistakes. Please verify important information before relying on it."
        )
        disclaimer_frame.paragraphs[0].font.size = Pt(10)
        disclaimer_frame.paragraphs[0].font.color.rgb = MUTED
        disclaimer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _bullet_slide(self, title: str, bullets: list[str]) -> None:
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        heading = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.7))
        heading.text_frame.text = title
        heading.text_frame.paragraphs[0].font.size = Pt(28)
        heading.text_frame.paragraphs[0].font.bold = True
        heading.text_frame.paragraphs[0].font.color.rgb = BLUE

        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2))
        text_frame = body.text_frame
        text_frame.clear()

        for index, bullet in enumerate(bullets[:8]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = SLATE

    def _chart_slides(self, chart_data: dict) -> None:
        chart_export = get_export_chart_images(chart_data)

        for title, png_bytes in chart_export.images:
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
            heading = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.7))
            heading.text_frame.text = title
            heading.text_frame.paragraphs[0].font.size = Pt(28)
            heading.text_frame.paragraphs[0].font.bold = True
            heading.text_frame.paragraphs[0].font.color.rgb = BLUE
            slide.shapes.add_picture(
                BytesIO(png_bytes), Inches(1.4), Inches(1.5), width=Inches(10.5)
            )

    def build(self, report_text: str) -> bytes:
        parsed = parse_intelligence_report(
            report_text,
            source_documents=self.metadata.source_documents,
            pack_type="board_pack",
        )
        metrics = dashboard_metrics(parsed)

        self._title_slide()

        dashboard_bullets = [
            f"{label}: {metrics[key]}{suffix}"
            for label, key, suffix in [
                ("Health Score", "health_score", "/100"),
                ("Outlook", "outlook", ""),
                ("Confidence", "confidence", ""),
                ("Documents analyzed", "documents", ""),
                ("Critical risks", "key_risks", ""),
                ("Recommendations", "recommendations", ""),
            ]
            if is_available(metrics.get(key))
        ]
        if dashboard_bullets:
            self._bullet_slide("Executive Dashboard", dashboard_bullets)

        self._chart_slides(parsed.chart_data)

        risks = top_risks(parsed)

        if risks:
            self._bullet_slide(
                "⚠ Key Risks",
                [f"{card['title']} — {card['severity']}" for card in risks],
            )

        recommendations = next(
            (
                section
                for section in parsed.sections
                if "recommendation" in section.title.lower()
            ),
            None,
        )

        if recommendations:
            # Recommendations are numbered list items whose Action clause
            # starts the line ("1. **Action:** ..."); Rationale/Measurement
            # are indented continuation lines with no leading marker, so
            # this pattern naturally picks up only the headline action.
            bullets = [
                re.sub(r"^\s*(?:[-*]|\d+[.)])\s*", "", line).replace("**", "").strip()
                for line in recommendations.body.splitlines()
                if re.match(r"^\s*(?:[-*]|\d+[.)])\s+", line)
            ]
            self._bullet_slide("💡 Recommendations", bullets or [recommendations.body[:300]])

        insight = first_ai_insight(parsed)

        if insight:
            self._bullet_slide("✦ AI Insights", [insight])

        trends = next(
            (section for section in parsed.sections if section.title.lower() == "trends"),
            None,
        )

        if trends:
            bullets = [
                line.strip("- ").strip()
                for line in trends.body.splitlines()
                if line.strip().startswith("-")
            ]
            self._bullet_slide("Trends", bullets or [trends.body[:300]])

        buffer = BytesIO()
        self.presentation.save(buffer)
        return buffer.getvalue()


def build_premium_presentation(
    *,
    report_text: str,
    metadata: PresentationExportMetadata,
) -> bytes:
    return PremiumPresentationBuilder(metadata).build(report_text)
