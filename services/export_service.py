"""
DataDumpAI Enterprise
Export Service

Central entry point for all report export operations.
"""

from __future__ import annotations

import re
import unicodedata
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    HRFlowable,
    Image,
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from services.export_chart_blocks import get_export_chart_images
from models.report_data import ReportData
from services.report_chart_data import prepare_report_for_output
from services.report_document import prepare_report_view, report_data_from_markdown
from services.visualization_engine import dashboard_section_heading
from services.report_markdown_renderer import (
    classify_label_value,
    drop_duplicate_leading_heading,
    highlight_value_html,
    parse_markdown_blocks,
    strip_inline_markdown,
)

# Shared brand palette — matches report_markdown_renderer.highlight_value_html
# and the existing table styling below, so PDF/DOCX/PPTX stay consistent.
BRAND_NAVY = "#0F172A"
BRAND_BLUE = "#1D4ED8"
BRAND_MUTED = "#64748B"
BRAND_RULE = "#93C5FD"
from storage.file_store import FileStore
from core.project_access import assert_project_access


class ExportService:
    """
    Handles report and project export operations.

    Every export format flows through this service so persistence,
    naming, and future format support stay in one place.
    """

    def __init__(self, *, access_token: str | None = None) -> None:
        self._access_token = access_token

    def _file_store(self) -> FileStore:
        return FileStore.for_current_user(access_token=self._access_token)

    MIME_TYPES = {
        "markdown": "text/markdown",
        "pdf": "application/pdf",
        "docx": (
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.document"
        ),
        "pptx": (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
    }

    def get_exports(self, project_id: str) -> list[dict[str, Any]]:
        try:
            assert_project_access(project_id)
        except PermissionError:
            return []

        store = self._file_store()
        exports: list[dict[str, Any]] = []

        for filename in store.list_files(project_id, "exports"):
            suffix = Path(filename).suffix.lower().lstrip(".")

            if suffix == "md":
                mime_type = self.MIME_TYPES["markdown"]
            elif suffix == "pdf":
                mime_type = self.MIME_TYPES["pdf"]
            elif suffix == "docx":
                mime_type = self.MIME_TYPES["docx"]
            else:
                mime_type = "application/octet-stream"

            if store._backend == "local":
                storage_path = str(store._local_root(project_id) / "exports" / filename)
            else:
                storage_path = store._storage_key(project_id, "exports", filename)

            try:
                size = len(store.read_bytes(storage_path))
            except Exception:
                size = 0

            exports.append(
                {
                    "filename": filename,
                    "path": storage_path,
                    "size": size,
                    "mime_type": mime_type,
                    "format": suffix or "unknown",
                    "exported_at": datetime.now(timezone.utc).isoformat(),
                }
            )

        return exports

    def _slugify(self, report_name: str) -> str:
        # Storage keys must be ASCII-safe — drop accents/em-dashes/curly
        # quotes/etc. (NFKD + ascii-ignore) rather than just denylisting a
        # few filesystem-unsafe characters, which let stray unicode like
        # "—" through and made Supabase Storage reject the upload outright.
        normalized = (
            unicodedata.normalize("NFKD", report_name)
            .encode("ascii", "ignore")
            .decode("ascii")
        )
        slug = re.sub(r"[^A-Za-z0-9]+", "_", normalized).strip("_").lower()
        return slug or "report"

    def _append_pdf_charts(self, story: list[Any], charts: list[tuple[str, bytes]], body_style: ParagraphStyle) -> None:
        story.append(Spacer(1, 0.12 * inch))

        for title, png_bytes in charts:
            story.append(Paragraph(f"<b>{title}</b>", body_style))
            story.append(
                Image(
                    BytesIO(png_bytes),
                    width=6.2 * inch,
                    height=2.6 * inch,
                    kind="proportional",
                )
            )
            story.append(Spacer(1, 0.15 * inch))

    def _append_docx_charts(self, document: Document, charts: list[tuple[str, bytes]]) -> None:
        for title, png_bytes in charts:
            heading = document.add_heading(title, level=3)
            for run in heading.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x1D, 0x4E, 0xD8)

            paragraph = document.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            paragraph.add_run().add_picture(BytesIO(png_bytes), width=Inches(6.2))

    def _prepend_pdf_charts(
        self,
        story: list[Any],
        chart_data: dict[str, Any],
        *,
        body_style: ParagraphStyle,
        heading_style: ParagraphStyle,
    ) -> None:
        chart_export = get_export_chart_images(chart_data)

        if not chart_export.images and not chart_export.unavailable_note:
            return

        story.append(Paragraph(dashboard_section_heading(chart_data), heading_style))

        if chart_export.images:
            self._append_pdf_charts(story, chart_export.images, body_style)

        if chart_export.unavailable_note:
            story.append(Paragraph(chart_export.unavailable_note, body_style))

    def _prepend_docx_charts(self, document: Document, chart_data: dict[str, Any]) -> None:
        chart_export = get_export_chart_images(chart_data)

        if not chart_export.images and not chart_export.unavailable_note:
            return

        document.add_heading(dashboard_section_heading(chart_data), level=2)

        if chart_export.images:
            self._append_docx_charts(document, chart_export.images)

        if chart_export.unavailable_note:
            paragraph = document.add_paragraph(chart_export.unavailable_note)
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    @staticmethod
    def _escape_pdf_text(text: str) -> str:
        return (
            strip_inline_markdown(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    @staticmethod
    def _shade_docx_cell(cell: Any, fill_hex: str) -> None:
        fill = fill_hex.lstrip("#").upper()
        tc_pr = cell._tc.get_or_add_tcPr()
        existing = tc_pr.find(qn("w:shd"))
        if existing is not None:
            tc_pr.remove(existing)
        shading = OxmlElement("w:shd")
        shading.set(qn("w:fill"), fill)
        shading.set(qn("w:val"), "clear")
        tc_pr.append(shading)

    def _render_pdf_table(
        self,
        rows: list[list[str]],
        *,
        story: list[Any],
        body_style: ParagraphStyle,
    ) -> None:
        if not rows:
            return

        col_count = max(len(row) for row in rows)
        col_width = 6.2 * inch / max(col_count, 1)
        cell_style = ParagraphStyle(
            "ExportTableCell",
            parent=body_style,
            fontSize=9,
            leading=12,
            alignment=TA_LEFT,
            spaceBefore=0,
            spaceAfter=0,
        )
        header_style = ParagraphStyle(
            "ExportTableHeader",
            parent=cell_style,
            textColor=colors.white,
        )

        table_data: list[list[Any]] = []
        for row_index, row in enumerate(rows):
            style = header_style if row_index == 0 else cell_style
            padded = list(row) + [""] * (col_count - len(row))
            table_data.append(
                [Paragraph(self._escape_pdf_text(cell), style) for cell in padded[:col_count]]
            )

        table = Table(table_data, colWidths=[col_width] * col_count, repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.11, 0.31, 0.85)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.Color(0.97, 0.98, 0.99)]),
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.Color(0.78, 0.83, 0.90)),
                    ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.Color(0.78, 0.83, 0.90)),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]
            )
        )
        story.append(table)
        story.append(Spacer(1, 0.12 * inch))

    def _render_docx_table(self, document: Document, rows: list[list[str]]) -> None:
        if not rows:
            return

        col_count = max(len(row) for row in rows)
        table = document.add_table(rows=len(rows), cols=col_count)
        table.style = "Table Grid"

        for row_index, row in enumerate(rows):
            padded = list(row) + [""] * (col_count - len(row))
            for col_index in range(col_count):
                cell = table.rows[row_index].cells[col_index]
                cell.text = strip_inline_markdown(padded[col_index])

                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

                if row_index == 0:
                    self._shade_docx_cell(cell, "1D4ED8")
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif row_index % 2 == 0:
                    self._shade_docx_cell(cell, "F8FAFC")

        document.add_paragraph()

    def _render_pdf_block(
        self,
        block: Any,
        *,
        story: list[Any],
        body_style: ParagraphStyle,
        heading_styles: dict[int, ParagraphStyle],
        list_item_style: ParagraphStyle,
    ) -> None:
        if block.block_type == "heading":
            style = heading_styles.get(block.level) or heading_styles[max(heading_styles)]
            story.append(Paragraph(strip_inline_markdown(block.content), style))
        elif block.block_type == "paragraph":
            story.append(Paragraph(self._escape_pdf_text(block.content), body_style))
        elif block.block_type == "bullets":
            story.append(
                ListFlowable(
                    [
                        ListItem(Paragraph(self._escape_pdf_text(item), list_item_style))
                        for item in block.items
                    ],
                    bulletType="bullet",
                    bulletColor=colors.HexColor(BRAND_BLUE),
                    bulletFontSize=8,
                    leftIndent=18,
                    spaceBefore=2,
                    spaceAfter=8,
                )
            )
        elif block.block_type == "numbered":
            story.append(
                ListFlowable(
                    [
                        ListItem(Paragraph(self._escape_pdf_text(item), list_item_style))
                        for item in block.items
                    ],
                    bulletType="1",
                    bulletFormat="%s.",
                    bulletColor=colors.HexColor(BRAND_NAVY),
                    bulletFontSize=10,
                    leftIndent=20,
                    spaceBefore=2,
                    spaceAfter=8,
                )
            )
        elif block.block_type == "label_value":
            story.append(
                Paragraph(highlight_value_html(block.label, block.value), body_style)
            )
        elif block.block_type == "table":
            self._render_pdf_table(block.rows, story=story, body_style=body_style)
        elif block.block_type == "spacer":
            story.append(Spacer(1, 0.08 * inch))

    def _render_docx_block(self, document: Document, block: Any) -> None:
        if block.block_type == "heading":
            document.add_heading(strip_inline_markdown(block.content), level=min(block.level, 4))
        elif block.block_type == "paragraph":
            paragraph = document.add_paragraph(strip_inline_markdown(block.content))
            paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
            paragraph.paragraph_format.line_spacing = 1.45
            paragraph.paragraph_format.space_before = Pt(6)
            paragraph.paragraph_format.space_after = Pt(10)
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        elif block.block_type == "bullets":
            for item in block.items:
                paragraph = document.add_paragraph(strip_inline_markdown(item), style="List Bullet")
                paragraph.paragraph_format.left_indent = Inches(0.35)
        elif block.block_type == "numbered":
            for item in block.items:
                paragraph = document.add_paragraph(strip_inline_markdown(item), style="List Number")
                paragraph.paragraph_format.left_indent = Inches(0.35)
        elif block.block_type == "label_value":
            label, value, color = classify_label_value(block.label, block.value)
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.space_before = Pt(2)
            paragraph.paragraph_format.space_after = Pt(8)
            label_run = paragraph.add_run(f"{label}: " if value else label)
            label_run.bold = True
            label_run.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
            if value:
                value_run = paragraph.add_run(value)
                value_run.bold = True
                value_run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        elif block.block_type == "table":
            self._render_docx_table(document, block.rows)

    def _build_result(
        self,
        *,
        project_id: str,
        filename: str,
        data: bytes,
        mime_type: str,
        report_name: str | None = None,
    ) -> dict[str, Any]:
        storage_path = self._file_store().write(project_id, "exports", filename, data)

        try:
            from services.activity_service import ActivityService

            export_label = (report_name or filename).strip()
            export_format = Path(filename).suffix.lstrip(".").upper() or "FILE"
            ActivityService(access_token=self._access_token).log(
                "export.downloaded",
                f"Downloaded {export_label} ({export_format})",
                metadata={"project_id": project_id, "filename": filename},
            )
        except Exception:
            pass

        return {
            "filename": filename,
            "path": storage_path,
            "data": data,
            "mime_type": mime_type,
            "size": len(data),
        }

    def _resolve_report(self, report: ReportData | str) -> ReportData:
        if isinstance(report, ReportData):
            return report
        return report_data_from_markdown(report)

    def _prepare_export(self, report: ReportData | str):
        return prepare_report_view(self._resolve_report(report))

    def _resolve_export_input(
        self,
        *,
        report: ReportData | str | None = None,
        report_text: str | None = None,
    ) -> ReportData | str:
        if report is not None:
            return report
        if report_text is not None:
            return report_text
        raise TypeError("export requires report or report_text")

    def export_markdown(
        self,
        *,
        project_id: str,
        report_name: str,
        report: ReportData | str | None = None,
        report_text: str | None = None,
    ) -> dict[str, Any]:
        """Export a report as Markdown."""

        filename = f"{self._slugify(report_name)}.md"
        data = self._prepare_export(
            self._resolve_export_input(report=report, report_text=report_text)
        ).text.encode("utf-8")

        return self._build_result(
            project_id=project_id,
            filename=filename,
            data=data,
            mime_type=self.MIME_TYPES["markdown"],
            report_name=report_name,
        )

    def _draw_pdf_footer(self, canvas: Any, doc: Any, *, report_name: str) -> None:
        canvas.saveState()
        y = 0.55 * inch
        canvas.setStrokeColor(colors.HexColor(BRAND_RULE))
        canvas.setLineWidth(0.6)
        canvas.line(0.9 * inch, y, letter[0] - 0.9 * inch, y)
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.HexColor(BRAND_MUTED))
        canvas.drawString(0.9 * inch, y - 14, strip_inline_markdown(report_name)[:90])
        canvas.drawRightString(letter[0] - 0.9 * inch, y - 14, f"Page {doc.page}")
        canvas.restoreState()

    def export_pdf(
        self,
        *,
        project_id: str,
        report_name: str,
        report: ReportData | str | None = None,
        report_text: str | None = None,
        workspace_name: str | None = None,
        period_name: str | None = None,
    ) -> dict[str, Any]:
        """Export a report as PDF."""

        prepared = self._prepare_export(
            self._resolve_export_input(report=report, report_text=report_text)
        )
        report_text = prepared.text
        filename = f"{self._slugify(report_name)}.pdf"
        buffer = BytesIO()

        document = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            leftMargin=0.9 * inch,
            rightMargin=0.9 * inch,
            topMargin=0.9 * inch,
            bottomMargin=0.9 * inch,
        )
        styles = getSampleStyleSheet()
        body_style = ParagraphStyle(
            "ExportBody",
            parent=styles["Normal"],
            fontSize=11,
            leading=16,
            alignment=TA_JUSTIFY,
            spaceBefore=6,
            spaceAfter=10,
            textColor=colors.HexColor(BRAND_NAVY),
        )
        list_item_style = ParagraphStyle(
            "ExportListItem",
            parent=body_style,
            alignment=TA_LEFT,
            spaceBefore=0,
            spaceAfter=4,
        )
        heading_styles = {
            2: ParagraphStyle(
                "ExportH2",
                parent=styles["Heading1"],
                fontSize=16,
                leading=20,
                spaceBefore=18,
                spaceAfter=10,
                textColor=colors.HexColor(BRAND_NAVY),
            ),
            3: ParagraphStyle(
                "ExportH3",
                parent=styles["Heading2"],
                fontSize=13,
                leading=17,
                spaceBefore=12,
                spaceAfter=6,
                textColor=colors.HexColor(BRAND_BLUE),
            ),
            4: ParagraphStyle(
                "ExportH4",
                parent=styles["Heading3"],
                fontSize=11.5,
                leading=15,
                spaceBefore=10,
                spaceAfter=4,
                textColor=colors.HexColor(BRAND_MUTED),
            ),
        }
        title_style = ParagraphStyle(
            "ExportTitle",
            parent=styles["Title"],
            fontSize=22,
            leading=27,
            alignment=TA_LEFT,
            textColor=colors.HexColor(BRAND_NAVY),
            spaceAfter=6,
        )
        meta_style = ParagraphStyle(
            "ExportMeta",
            parent=styles["Normal"],
            fontSize=9.5,
            textColor=colors.HexColor(BRAND_MUTED),
            spaceBefore=0,
            spaceAfter=20,
        )

        story: list[Any] = [
            Paragraph(self._escape_pdf_text(report_name), title_style),
            HRFlowable(
                width="100%",
                thickness=1.4,
                color=colors.HexColor(BRAND_RULE),
                spaceAfter=6,
            ),
        ]
        meta_parts = [part for part in (workspace_name, period_name) if part]
        meta_parts.append(f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y')}")
        story.append(
            Paragraph(
                " &nbsp;&middot;&nbsp; ".join(self._escape_pdf_text(part) for part in meta_parts),
                meta_style,
            )
        )

        self._prepend_pdf_charts(
            story,
            prepared.chart_data,
            body_style=body_style,
            heading_style=heading_styles[2],
        )

        blocks = drop_duplicate_leading_heading(
            parse_markdown_blocks(report_text), report_name
        )
        for block in blocks:
            self._render_pdf_block(
                block,
                story=story,
                body_style=body_style,
                heading_styles=heading_styles,
                list_item_style=list_item_style,
            )

        document.build(
            story,
            onFirstPage=lambda canvas, doc: self._draw_pdf_footer(
                canvas, doc, report_name=report_name
            ),
            onLaterPages=lambda canvas, doc: self._draw_pdf_footer(
                canvas, doc, report_name=report_name
            ),
        )
        data = buffer.getvalue()

        return self._build_result(
            project_id=project_id,
            filename=filename,
            data=data,
            mime_type=self.MIME_TYPES["pdf"],
            report_name=report_name,
        )

    @staticmethod
    def _add_docx_bottom_rule(paragraph: Any, *, color_hex: str = "93C5FD", size: int = 10) -> None:
        p_pr = paragraph._p.get_or_add_pPr()
        borders = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), str(size))
        bottom.set(qn("w:space"), "4")
        bottom.set(qn("w:color"), color_hex)
        borders.append(bottom)
        p_pr.append(borders)

    @staticmethod
    def _add_docx_page_number_field(
        paragraph: Any, *, color_hex: str = "64748B", size: int = 8
    ) -> None:
        run = paragraph.add_run()
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color_hex)
        begin = OxmlElement("w:fldChar")
        begin.set(qn("w:fldCharType"), "begin")
        instr = OxmlElement("w:instrText")
        instr.set(qn("xml:space"), "preserve")
        instr.text = "PAGE"
        end = OxmlElement("w:fldChar")
        end.set(qn("w:fldCharType"), "end")
        run._r.append(begin)
        run._r.append(instr)
        run._r.append(end)

    def _style_docx_headings(self, document: Document) -> None:
        palette = {1: BRAND_NAVY, 2: BRAND_NAVY, 3: BRAND_BLUE, 4: BRAND_MUTED}
        for level, hex_color in palette.items():
            try:
                style = document.styles[f"Heading {level}"]
            except KeyError:
                continue
            style.font.color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
            style.font.bold = True

    def export_docx(
        self,
        *,
        project_id: str,
        report_name: str,
        report: ReportData | str | None = None,
        report_text: str | None = None,
        workspace_name: str | None = None,
        period_name: str | None = None,
    ) -> dict[str, Any]:
        """Export a report as Word (.docx)."""

        prepared = self._prepare_export(
            self._resolve_export_input(report=report, report_text=report_text)
        )
        report_text = prepared.text
        filename = f"{self._slugify(report_name)}.docx"
        buffer = BytesIO()

        document = Document()
        for section in document.sections:
            section.top_margin = Inches(0.9)
            section.bottom_margin = Inches(0.9)
            section.left_margin = Inches(0.9)
            section.right_margin = Inches(0.9)

        self._style_docx_headings(document)

        title_paragraph = document.add_heading(strip_inline_markdown(report_name), level=1)
        self._add_docx_bottom_rule(title_paragraph)

        meta_parts = [part for part in (workspace_name, period_name) if part]
        meta_parts.append(f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y')}")
        meta_paragraph = document.add_paragraph(" · ".join(meta_parts))
        meta_paragraph.paragraph_format.space_before = Pt(4)
        meta_paragraph.paragraph_format.space_after = Pt(18)
        for run in meta_paragraph.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = RGBColor.from_string(BRAND_MUTED.lstrip("#"))

        self._prepend_docx_charts(document, prepared.chart_data)

        blocks = drop_duplicate_leading_heading(
            parse_markdown_blocks(report_text), report_name
        )
        for block in blocks:
            self._render_docx_block(document, block)

        footer_paragraph = document.sections[0].footer.paragraphs[0]
        footer_paragraph.text = strip_inline_markdown(report_name)[:90] + "   "
        for run in footer_paragraph.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor.from_string(BRAND_MUTED.lstrip("#"))
        page_label_run = footer_paragraph.add_run("Page ")
        page_label_run.font.size = Pt(8)
        page_label_run.font.color.rgb = RGBColor.from_string(BRAND_MUTED.lstrip("#"))
        self._add_docx_page_number_field(footer_paragraph)

        document.save(buffer)
        data = buffer.getvalue()

        return self._build_result(
            project_id=project_id,
            filename=filename,
            data=data,
            mime_type=self.MIME_TYPES["docx"],
            report_name=report_name,
        )

    @staticmethod
    def _group_pptx_sections(blocks: list[Any]) -> list[tuple[str, list[Any]]]:
        """Split a flat block list into (H2 title, child blocks) sections —
        one slide group per top-level section, mirroring how the document
        body reads visually in the PDF/DOCX exports."""

        sections: list[tuple[str, list[Any]]] = []
        current_title = "Overview"
        current_blocks: list[Any] = []
        for block in blocks:
            if block.block_type == "heading" and block.level == 2:
                if current_blocks:
                    sections.append((current_title, current_blocks))
                current_title = strip_inline_markdown(block.content)
                current_blocks = []
                continue
            current_blocks.append(block)
        if current_blocks:
            sections.append((current_title, current_blocks))
        return sections

    @staticmethod
    def _set_pptx_bullet(paragraph: Any, *, color_hex: str) -> None:
        from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement

        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("marL", "228600")
        p_pr.set("indent", "-228600")
        bu_clr = PptxOxmlElement("a:buClr")
        srgb = PptxOxmlElement("a:srgbClr")
        srgb.set("val", color_hex.lstrip("#").upper())
        bu_clr.append(srgb)
        bu_font = PptxOxmlElement("a:buFont")
        bu_font.set("typeface", "Arial")
        bu_char = PptxOxmlElement("a:buChar")
        bu_char.set("char", "•")
        p_pr.append(bu_clr)
        p_pr.append(bu_font)
        p_pr.append(bu_char)

    @staticmethod
    def _set_pptx_number(paragraph: Any, *, start_at: int, color_hex: str) -> None:
        from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement

        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("marL", "274638")
        p_pr.set("indent", "-274638")
        bu_clr = PptxOxmlElement("a:buClr")
        srgb = PptxOxmlElement("a:srgbClr")
        srgb.set("val", color_hex.lstrip("#").upper())
        bu_clr.append(srgb)
        bu_font = PptxOxmlElement("a:buFont")
        bu_font.set("typeface", "Arial")
        bu_auto_num = PptxOxmlElement("a:buAutoNum")
        bu_auto_num.set("type", "arabicPeriod")
        bu_auto_num.set("startAt", str(max(1, start_at)))
        p_pr.append(bu_clr)
        p_pr.append(bu_font)
        p_pr.append(bu_auto_num)

    def _add_pptx_section_slides(
        self,
        presentation: Any,
        layout: Any,
        title: str,
        blocks: list[Any],
    ) -> None:
        from pptx.dml.color import RGBColor as PptxRGBColor
        from pptx.util import Inches as PptxInches, Pt as PptxPt

        navy = PptxRGBColor.from_string(BRAND_NAVY.lstrip("#"))
        blue = PptxRGBColor.from_string(BRAND_BLUE.lstrip("#"))

        # Flatten blocks into renderable (text, kind, meta) lines — tables
        # are rare in these narrative reports, so they render as plain
        # pipe-separated rows rather than a real pptx table (out of scope).
        lines: list[tuple[str, str, Any]] = []
        for block in blocks:
            if block.block_type == "heading":
                lines.append((block.content, "subheading", None))
            elif block.block_type == "paragraph":
                lines.append((block.content, "paragraph", None))
            elif block.block_type == "bullets":
                for item in block.items:
                    lines.append((item, "bullet", None))
            elif block.block_type == "numbered":
                for index, item in enumerate(block.items, start=1):
                    lines.append((item, "numbered", index))
            elif block.block_type == "label_value":
                label, value, color = classify_label_value(block.label, block.value)
                text = f"{label}: {value}" if value else label
                lines.append((text, "label_value", color))
            elif block.block_type == "table":
                for row in block.rows:
                    lines.append((" · ".join(row), "paragraph", None))

        if not lines:
            return

        chunk_size = 7
        chunks = [lines[i : i + chunk_size] for i in range(0, len(lines), chunk_size)]

        for chunk_index, chunk in enumerate(chunks):
            slide = presentation.slides.add_slide(layout)

            heading_box = slide.shapes.add_textbox(
                PptxInches(0.7), PptxInches(0.5), PptxInches(12.0), PptxInches(0.9)
            )
            heading_text = title if chunk_index == 0 else f"{title} (cont'd)"
            heading_frame = heading_box.text_frame
            heading_frame.word_wrap = True
            heading_frame.text = strip_inline_markdown(heading_text)
            heading_frame.paragraphs[0].font.size = PptxPt(28)
            heading_frame.paragraphs[0].font.bold = True
            heading_frame.paragraphs[0].font.color.rgb = navy

            body_box = slide.shapes.add_textbox(
                PptxInches(0.9), PptxInches(1.55), PptxInches(11.5), PptxInches(5.5)
            )
            text_frame = body_box.text_frame
            text_frame.word_wrap = True

            for item_index, (text, kind, meta) in enumerate(chunk):
                paragraph = text_frame.paragraphs[0] if item_index == 0 else text_frame.add_paragraph()
                clean_text = strip_inline_markdown(text)
                paragraph.text = clean_text
                paragraph.space_after = PptxPt(10)

                if kind == "subheading":
                    paragraph.font.size = PptxPt(18)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = blue
                elif kind == "bullet":
                    paragraph.font.size = PptxPt(16)
                    paragraph.font.color.rgb = navy
                    self._set_pptx_bullet(paragraph, color_hex=BRAND_BLUE)
                elif kind == "numbered":
                    paragraph.font.size = PptxPt(16)
                    paragraph.font.color.rgb = navy
                    self._set_pptx_number(paragraph, start_at=meta, color_hex=BRAND_NAVY)
                elif kind == "label_value":
                    paragraph.font.size = PptxPt(14)
                    paragraph.font.italic = True
                    paragraph.font.color.rgb = PptxRGBColor.from_string(
                        str(meta).lstrip("#")
                    )
                else:
                    paragraph.font.size = PptxPt(15)
                    paragraph.font.color.rgb = navy

    def export_pptx(
        self,
        *,
        project_id: str,
        report_name: str,
        report: ReportData | str | None = None,
        report_text: str | None = None,
        workspace_name: str | None = None,
        period_name: str | None = None,
    ) -> dict[str, Any]:
        """Export a report as PowerPoint (.pptx).

        Generic, markdown-block-driven builder for report types that aren't
        the specialized "Executive Intelligence Dashboard" shape handled by
        premium_pptx_export.py — reuses that module's visual palette/sizing
        but drives slide content from parse_markdown_blocks() output.
        """

        from pptx import Presentation
        from pptx.dml.color import RGBColor as PptxRGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches as PptxInches, Pt as PptxPt

        prepared = self._prepare_export(
            self._resolve_export_input(report=report, report_text=report_text)
        )
        report_text = prepared.text
        filename = f"{self._slugify(report_name)}.pptx"

        presentation = Presentation()
        presentation.slide_width = PptxInches(13.33)
        presentation.slide_height = PptxInches(7.5)
        blank_layout = presentation.slide_layouts[6]

        navy = PptxRGBColor.from_string(BRAND_NAVY.lstrip("#"))
        blue = PptxRGBColor.from_string(BRAND_BLUE.lstrip("#"))

        title_slide = presentation.slides.add_slide(blank_layout)
        title_box = title_slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(2.4), PptxInches(11.7), PptxInches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = strip_inline_markdown(report_name)
        title_frame.paragraphs[0].font.size = PptxPt(34)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = navy
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        meta_parts = [part for part in (workspace_name, period_name) if part]
        meta_parts.append(f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y')}")
        meta_box = title_slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(3.95), PptxInches(11.7), PptxInches(0.8)
        )
        meta_frame = meta_box.text_frame
        meta_frame.text = "  ·  ".join(meta_parts)
        meta_frame.paragraphs[0].font.size = PptxPt(16)
        meta_frame.paragraphs[0].font.color.rgb = blue
        meta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Unlike PDF/DOCX (where the body sits directly under a rendered
        # title and an immediately-repeated heading looks broken), each
        # PPTX section gets its own slide — the leading "Executive Summary"
        # heading is a legitimate, wanted slide title here, so it is not
        # dropped even though it overlaps the report title.
        blocks = parse_markdown_blocks(report_text)
        for section_title, section_blocks in self._group_pptx_sections(blocks):
            self._add_pptx_section_slides(presentation, blank_layout, section_title, section_blocks)

        buffer = BytesIO()
        presentation.save(buffer)
        data = buffer.getvalue()

        return self._build_result(
            project_id=project_id,
            filename=filename,
            data=data,
            mime_type=self.MIME_TYPES["pptx"],
            report_name=report_name,
        )
