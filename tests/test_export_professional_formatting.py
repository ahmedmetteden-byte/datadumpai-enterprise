"""Smoke tests for the professional PDF/DOCX/PPTX export redesign: real
numbered/bulleted lists (not run-on paragraphs or baked checkmarks), no
duplicated title heading, and a working generic PPTX path."""

from __future__ import annotations

from io import BytesIO

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

from services.export_service import ExportService

REPORT_WITH_LISTS = """
## Executive Summary

This is a concise summary of the situation.

## Risks & Issues

- Regulatory capital requirements may tighten in the next cycle
- Reinsurance costs are trending upward

## Strategic Recommendations

1. Commission an independent reserve adequacy review this quarter.
2. Renegotiate reinsurance treaty terms ahead of renewal.
3. Launch a cross-sell pilot into the commercial lines segment.
"""


def test_export_pdf_renders_numbered_list_items_separately():
    result = ExportService().export_pdf(
        project_id="export-format-project",
        report_name="Executive Summary — Custom / Ad hoc",
        report_text=REPORT_WITH_LISTS,
    )

    assert result["data"].startswith(b"%PDF")

    pdf_text = "\n".join(
        page.extract_text() or "" for page in PdfReader(BytesIO(result["data"])).pages
    )

    assert "Commission an independent reserve adequacy review" in pdf_text
    assert "Renegotiate reinsurance treaty terms" in pdf_text
    assert "✓" not in pdf_text
    # The title/footer legitimately repeat the report name, but the body's
    # own leading "## Executive Summary" heading must be dropped rather
    # than rendered a second time right under the title block — the
    # metadata line should be followed directly by the summary paragraph.
    assert "Generated" in pdf_text
    metadata_index = pdf_text.index("Generated")
    assert "This is a concise summary" in pdf_text[metadata_index:metadata_index + 120]


def test_export_docx_renders_numbered_and_bullet_styles():
    result = ExportService().export_docx(
        project_id="export-format-project",
        report_name="Executive Summary — Custom / Ad hoc",
        report_text=REPORT_WITH_LISTS,
    )

    assert result["data"].startswith(b"PK")

    document = Document(BytesIO(result["data"]))
    styles_used = {p.style.name for p in document.paragraphs if p.text.strip()}

    assert "List Number" in styles_used
    assert "List Bullet" in styles_used

    non_empty_paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    assert "✓" not in "\n".join(non_empty_paragraphs)
    # First heading is the title; the body must not repeat "Executive
    # Summary" as its own duplicate heading right after the metadata line.
    heading_texts = [
        p.text for p in document.paragraphs if p.style.name.startswith("Heading")
    ]
    assert heading_texts.count("Executive Summary") == 0


def test_export_pptx_builds_title_slide_and_section_slides():
    result = ExportService().export_pptx(
        project_id="export-format-project",
        report_name="Executive Summary — Custom / Ad hoc",
        report_text=REPORT_WITH_LISTS,
        workspace_name="Acme Insurance Co",
        period_name="Custom / Ad hoc",
    )

    assert result["filename"].endswith(".pptx")
    assert result["data"].startswith(b"PK")

    presentation = Presentation(BytesIO(result["data"]))
    slides = list(presentation.slides)
    assert len(slides) >= 3

    all_text = []
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)

    assert "Executive Summary" in joined
    assert "Risks & Issues" in joined
    assert "Strategic Recommendations" in joined
    assert "Commission an independent reserve adequacy review" in joined
    assert "✓" not in joined
