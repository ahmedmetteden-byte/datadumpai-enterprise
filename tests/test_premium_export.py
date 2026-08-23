"""Tests for premium PDF and presentation export."""

from __future__ import annotations

from services.premium_docx_export import DocxExportMetadata, build_premium_docx
from services.premium_pdf_export import PremiumExportMetadata, build_premium_pdf
from services.premium_pptx_export import PresentationExportMetadata, build_premium_presentation
from services.report_document_parser import parse_intelligence_report

SAMPLE_REPORT = """
## Executive Intelligence Dashboard

### Executive Summary Card
| Field | Value |
| --- | --- |
| Industry Status | 🟡 Cautious |
| Confidence | 90% |
| Priority | Claims Reform |
| Overall Trend | Improving |

### Executive Snapshot
| Metric | Value |
| --- | --- |
| Documents analyzed | 4 |
| Reporting period | December 2022 – November 2024 |
| Critical risks | 4 |
| Recommendations | 7 |
| Overall outlook | Cautious |
| AI confidence | 90% |

### Top Risks
- 🔴 **Claims** — Settlement delays continue
- 🟠 **Capital** — Adequacy pressure remains

## AI Insights
- Claims issues increased across all four meetings.

<!-- REPORT_CHARTS
{
  "topics": [{"label": "Claims", "value": 31}, {"label": "Capital", "value": 21}],
  "health_score": 75
}
-->
"""


def test_parse_intelligence_report_extracts_card_and_charts():
    parsed = parse_intelligence_report(
        SAMPLE_REPORT,
        source_documents=["meeting.pdf", "report.pdf"],
    )

    assert parsed.summary_card["Priority"] == "Claims Reform"
    assert parsed.chart_data["health_score"] == 75
    assert len(parsed.source_documents) == 2


def test_build_premium_pdf_returns_pdf_bytes():
    pdf_bytes = build_premium_pdf(
        report_text=SAMPLE_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Board Meeting - November 2024",
            report_name="Executive Summary",
            reporting_period="December 2022 – November 2024",
            source_documents=["meeting.pdf"],
            pack_type="executive",
        ),
    )

    assert pdf_bytes.startswith(b"%PDF")


class _FakeCanvas:
    """Records the draw calls _draw_page_decorations makes, without
    depending on real PDF rendering/text-extraction (rotated watermark
    text has been shown, elsewhere in this codebase, to extract unreliably
    via naive PDF-text tools -- testing the actual method calls the
    watermark logic branches on is more direct and robust)."""

    def __init__(self) -> None:
        self.calls: list[str] = []

    def saveState(self):
        self.calls.append("saveState")

    def restoreState(self):
        self.calls.append("restoreState")

    def setFillColor(self, *_args, **_kwargs):
        self.calls.append("setFillColor")

    def setFont(self, *_args, **_kwargs):
        self.calls.append("setFont")

    def translate(self, *_args, **_kwargs):
        self.calls.append("translate")

    def rotate(self, *_args, **_kwargs):
        self.calls.append("rotate")

    def drawCentredString(self, _x, _y, text):
        self.calls.append(f"drawCentredString:{text}")

    def setStrokeColor(self, *_args, **_kwargs):
        pass

    def line(self, *_args, **_kwargs):
        pass

    def drawString(self, _x, _y, text):
        self.calls.append(f"drawString:{text}")

    def drawRightString(self, _x, _y, text):
        self.calls.append(f"drawRightString:{text}")


class _FakeDoc:
    def __init__(self, page: int) -> None:
        self.page = page


def test_free_plan_pdf_draws_the_watermark():
    from services.premium_pdf_export import PremiumPDFBuilder

    builder = PremiumPDFBuilder(
        PremiumExportMetadata(
            project_name="p", report_name="r", show_watermark=True,
        )
    )
    canvas = _FakeCanvas()
    builder._draw_page_decorations(canvas, _FakeDoc(page=1))

    assert "rotate" in canvas.calls
    assert "drawCentredString:DataDumpAI" in canvas.calls


def test_paid_plan_pdf_omits_the_watermark():
    from services.premium_pdf_export import PremiumPDFBuilder

    builder = PremiumPDFBuilder(
        PremiumExportMetadata(
            project_name="p", report_name="r", show_watermark=False,
        )
    )
    canvas = _FakeCanvas()
    builder._draw_page_decorations(canvas, _FakeDoc(page=1))

    assert "rotate" not in canvas.calls
    assert "drawCentredString:DataDumpAI" not in canvas.calls


def test_paid_plan_pdf_keeps_the_header_footer_branding_on_later_pages():
    """Removing the watermark must not remove the (separate, non-diagonal)
    header/footer branding that appears on page 2+."""

    from services.premium_pdf_export import PremiumPDFBuilder

    builder = PremiumPDFBuilder(
        PremiumExportMetadata(
            project_name="p", report_name="r", show_watermark=False,
        )
    )
    canvas = _FakeCanvas()
    builder._draw_page_decorations(canvas, _FakeDoc(page=2))

    assert "rotate" not in canvas.calls
    # The header/footer branding (a separate, non-diagonal element) still
    # renders even with the watermark off.
    assert "drawString:DataDumpAI" in canvas.calls
    assert any(call.startswith("drawRightString:Page") for call in canvas.calls)
    # Only the footer's saveState/restoreState pair ran -- the watermark's
    # own pair was skipped, not left unbalanced.
    assert canvas.calls.count("saveState") == 1
    assert canvas.calls.count("restoreState") == 1


def test_build_premium_pdf_returns_pdf_bytes():
    pptx_bytes = build_premium_presentation(
        report_text=SAMPLE_REPORT,
        metadata=PresentationExportMetadata(
            project_name="Board Meeting - November 2024",
            report_name="Executive Summary",
            source_documents=["meeting.pdf"],
        ),
    )

    assert pptx_bytes.startswith(b"PK")


# Step G: the premium renderer must also work for a real SPA-format
# report (no Executive Intelligence Dashboard heading, no health_score),
# not only the legacy intelligence format above.
SPA_REPORT = """
## Executive Summary
Gross premiums increased 97.4% from 2022 to 2024.

## Key Findings

### Gross Premium increased 97.4%
Detail text.
**Basis:** Calculated result
**Confidence:** High — verified calculations.
**Source:** report.xlsx

## Risks & Issues
- **Rising Claims Costs:** The significant increase in Gross Claims poses a risk to profitability.

## Opportunities
- **Enhanced Claims Management:** Refine claims management processes.

## Strategic Recommendations
1. **Action:** Conduct a comprehensive review of claims management processes.
   **Rationale:** The 51.3% increase in Gross Claims suggests inefficiencies.
   **Measurement:** Track the ratio of Gross Claims to Gross Premium.
2. **Action:** Explore new market segments.
   **Rationale:** The disparity in growth rates indicates untapped potential.
   **Measurement:** Evaluate new segments based on premium growth metrics.

## Conclusion
Momentum should be sustained.

<!-- REPORT_CHARTS
{"visualizations": [{"title": "Gross Premium", "type": "LINE_CHART", "description": "d", "data": {"trends": [{"label": "2024", "prior": 1043.1, "current": 1558.7}]}, "priority": 1, "decision_question": "", "unit": "", "x_label": "Period", "y_label": "Gross Premium"}]}
-->
"""


EVIDENCE_REPORT = """
## Executive Summary
Gross premiums increased 97.4% from 2022 to 2024.

## Key Findings

### Gross Premium increased 97.4%
Detail text about the finding.
**Basis:** Calculated result
**Confidence:** High — verified calculations.
**Source:** Annual-Statistical-Market-Report-Updated-01-2023.pdf, Annual-Statistical-Market-Report-2024.pdf

## Conclusion
Momentum should be sustained.
"""


def test_build_premium_pdf_renders_evidence_caption_and_humanized_source_filenames():
    """Report Output Quality Upgrade Step A: Basis/Confidence/Source must
    render as a visually subordinate "Evidence" block with human-readable
    source titles in the Key Finding itself. The real filename is still
    expected to appear verbatim in the separate Source References appendix
    (report_document_parser.py always lists real source_documents there) —
    only the inline evidence tag is humanized, nothing is fabricated or
    hidden."""

    import io

    from PyPDF2 import PdfReader

    pdf_bytes = build_premium_pdf(
        report_text=EVIDENCE_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=[
                "Annual-Statistical-Market-Report-Updated-01-2023.pdf",
                "Annual-Statistical-Market-Report-2024.pdf",
            ],
            pack_type="executive",
        ),
    )
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf_bytes)).pages)

    assert "EVIDENCE" in pdf_text
    assert "Annual Statistical Market Report 2023" in pdf_text
    assert "Annual Statistical Market Report 2024" in pdf_text

    finding_section = pdf_text.split("Gross Premium increased 97.4%", 1)[1].split("Conclusion", 1)[0]
    assert "Annual-Statistical-Market-Report-Updated-01-2023.pdf" not in finding_section


def test_build_premium_docx_renders_evidence_caption_and_humanized_source_filenames():
    docx_bytes = build_premium_docx(
        report_text=EVIDENCE_REPORT,
        metadata=DocxExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=[
                "Annual-Statistical-Market-Report-Updated-01-2023.pdf",
                "Annual-Statistical-Market-Report-2024.pdf",
            ],
            pack_type="executive",
        ),
    )

    from io import BytesIO

    from docx import Document as DocxDocument

    document = DocxDocument(BytesIO(docx_bytes))
    docx_text = "\n".join(p.text for p in document.paragraphs)

    assert "EVIDENCE" in docx_text
    assert "Annual Statistical Market Report 2023" in docx_text
    assert "Annual Statistical Market Report 2024" in docx_text

    finding_section = docx_text.split("Gross Premium increased 97.4%", 1)[1].split("Conclusion", 1)[0]
    assert "Annual-Statistical-Market-Report-Updated-01-2023.pdf" not in finding_section


NO_RISKS_OR_OPPORTUNITIES_REPORT = """
## Executive Summary
Gross premiums increased 97.4% from 2022 to 2024.

## Risks & Issues
No risks were identified in the evidence reviewed.

## Opportunities
No opportunities were identified in the evidence reviewed.

## Conclusion
Momentum should be sustained.
"""


def test_build_premium_pdf_softens_no_risks_no_opportunities_wording():
    """Report Output Quality Upgrade Step B: 'No critical risks identified.'
    overstates the claim — it reads as 'none exist' when the honest claim
    is 'none were found in the evidence reviewed'."""

    import io

    from PyPDF2 import PdfReader

    pdf_bytes = build_premium_pdf(
        report_text=NO_RISKS_OR_OPPORTUNITIES_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf_bytes)).pages)

    assert "No risks were identified in the evidence reviewed." in pdf_text
    assert "No opportunities were identified in the evidence reviewed." in pdf_text
    assert "No critical risks identified." not in pdf_text
    assert "No opportunities identified." not in pdf_text


def test_build_premium_docx_softens_no_risks_no_opportunities_wording():
    docx_bytes = build_premium_docx(
        report_text=NO_RISKS_OR_OPPORTUNITIES_REPORT,
        metadata=DocxExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )

    from io import BytesIO

    from docx import Document as DocxDocument

    docx_text = "\n".join(p.text for p in DocxDocument(BytesIO(docx_bytes)).paragraphs)

    assert "No risks were identified in the evidence reviewed." in docx_text
    assert "No opportunities were identified in the evidence reviewed." in docx_text
    assert "No critical risks identified." not in docx_text
    assert "No opportunities identified." not in docx_text


def test_executive_summary_page_keeps_top_risks_heading_with_its_content():
    """Report Output Quality Upgrade Step F: 'Top Opportunities' was seen
    landing alone at the bottom of a page in a real generated PDF, with
    its actual bullets starting fresh on the next page — no heading in
    sight. The heading and at least its first content item must be in the
    same KeepTogether flowable so ReportLab can never split them apart."""

    from reportlab.platypus import KeepTogether, Paragraph

    from services.premium_pdf_export import PremiumExportMetadata, PremiumPDFBuilder
    from services.report_document_parser import parse_intelligence_report

    parsed = parse_intelligence_report(SPA_REPORT, source_documents=["report.xlsx"])
    builder = PremiumPDFBuilder(
        PremiumExportMetadata(
            project_name="p", report_name="r", reporting_period="x",
            source_documents=["report.xlsx"], pack_type="executive",
        )
    )
    story = builder._executive_summary_page(parsed)

    keep_togethers = [item for item in story if isinstance(item, KeepTogether)]
    assert len(keep_togethers) >= 2, "expected separate KeepTogether blocks for risks and opportunities"

    def _flowable_text(flowable) -> str:
        frags = getattr(flowable, "frags", None)
        if frags:
            return "".join(f.text for f in frags)
        return ""

    def _heading_of(block) -> str:
        first = block._content[0]
        return _flowable_text(first) if isinstance(first, Paragraph) else ""

    # Indexed by heading text, not raw position — the chart section (also
    # correctly KeepTogether-wrapped, so a chart heading can never split
    # from its own image across a page break) may add its own blocks
    # ahead of Risks/Opportunities.
    by_heading = {_heading_of(block): block for block in keep_togethers}

    risks_block = by_heading["Top Risks"]
    assert isinstance(risks_block._content[0], Paragraph)
    assert len(risks_block._content) >= 2  # heading + at least one content line

    opportunities_block = by_heading["Top Opportunities"]
    assert len(opportunities_block._content) >= 2


def test_appendix_heading_stays_with_first_subsection():
    """The 'Appendix' section title used to be a standalone top-level
    flowable, so it could land alone at the bottom of a page with the
    first appendix subsection (e.g. Source References) starting fresh on
    the next page. It must be merged into the first subsection's
    KeepTogether block."""

    from reportlab.platypus import KeepTogether, PageBreak, Paragraph

    from services.premium_pdf_export import PremiumExportMetadata, PremiumPDFBuilder
    from services.report_document_parser import parse_intelligence_report

    parsed = parse_intelligence_report(SPA_REPORT, source_documents=["report.xlsx"])
    builder = PremiumPDFBuilder(
        PremiumExportMetadata(
            project_name="p", report_name="r", reporting_period="x",
            source_documents=["report.xlsx"], pack_type="executive",
        )
    )
    story = builder._appendix_story(parsed)

    assert isinstance(story[0], PageBreak)
    first_block = story[1]
    assert isinstance(first_block, KeepTogether)

    def _flowable_text(flowable) -> str:
        frags = getattr(flowable, "frags", None)
        return "".join(f.text for f in frags) if frags else ""

    assert isinstance(first_block._content[0], Paragraph)
    assert _flowable_text(first_block._content[0]) == "Appendix"
    assert _flowable_text(first_block._content[1]) == "Source References"


def test_docx_headings_are_kept_with_next_paragraph():
    """DOCX equivalent of the PDF KeepTogether fix — every heading gets
    Word's 'keep with next' formatting so a heading can never be the last
    thing on a page with its content starting fresh on the next."""

    from docx import Document as DocxDocument
    from io import BytesIO

    docx_bytes = build_premium_docx(
        report_text=SPA_REPORT,
        metadata=DocxExportMetadata(
            project_name="p", report_name="r", reporting_period="x",
            source_documents=["report.xlsx"], pack_type="executive",
        ),
    )
    document = DocxDocument(BytesIO(docx_bytes))

    headings = [p for p in document.paragraphs if p.style.name.startswith("Heading")]
    assert headings, "expected at least one heading paragraph"
    assert all(p.paragraph_format.keep_with_next for p in headings)


def test_build_premium_pdf_handles_real_spa_format_report_without_crashing():
    pdf_bytes = build_premium_pdf(
        report_text=SPA_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )
    assert pdf_bytes.startswith(b"%PDF")


def test_build_premium_docx_handles_real_spa_format_report_without_crashing():
    docx_bytes = build_premium_docx(
        report_text=SPA_REPORT,
        metadata=DocxExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )
    assert docx_bytes.startswith(b"PK")


def test_build_premium_pdf_renders_recommendation_action_clause_not_just_rationale():
    """Report Output Quality Upgrade Step C: a numbered recommendation's
    **Action:** clause was silently dropped by _render_block() (no
    "numbered" block-type case existed), so only Rationale/Measurement
    ever appeared — the recommendation's own title/action was missing."""

    import io

    from PyPDF2 import PdfReader

    pdf_bytes = build_premium_pdf(
        report_text=SPA_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf_bytes)).pages)

    assert "Conduct a comprehensive review of claims management processes" in pdf_text
    assert "Explore new market segments" in pdf_text
    assert "51.3% increase in Gross Claims suggests inefficiencies" in pdf_text


def test_build_premium_docx_renders_recommendation_action_clause_not_just_rationale():
    docx_bytes = build_premium_docx(
        report_text=SPA_REPORT,
        metadata=DocxExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )

    from io import BytesIO

    from docx import Document as DocxDocument

    docx_text = "\n".join(p.text for p in DocxDocument(BytesIO(docx_bytes)).paragraphs)

    assert "Conduct a comprehensive review of claims management processes" in docx_text
    assert "Explore new market segments" in docx_text


def test_build_premium_pdf_drops_isolated_singular_recommendation_preview():
    """The early one-sentence 'Strategic Recommendation' preview on the
    executive-summary page duplicated the full section with weaker text
    and reliably produced an isolated near-empty page — dropped rather
    than replaced, per Step C."""

    import io

    from PyPDF2 import PdfReader

    pdf_bytes = build_premium_pdf(
        report_text=SPA_REPORT,
        metadata=PremiumExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            reporting_period="Custom / Ad hoc",
            source_documents=["report.xlsx"],
            pack_type="executive",
        ),
    )
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf_bytes)).pages)

    assert "Strategic Recommendation\n" not in pdf_text
    # The full, plural section must still be present.
    assert "Strategic Recommendations" in pdf_text


def test_build_premium_presentation_embeds_charts_for_spa_format_report():
    pptx_bytes = build_premium_presentation(
        report_text=SPA_REPORT,
        metadata=PresentationExportMetadata(
            project_name="Q4 Revenue Review",
            report_name="Q4 Revenue Report",
            source_documents=["report.xlsx"],
        ),
    )
    assert pptx_bytes.startswith(b"PK")

    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(pptx_bytes))
    has_picture = any(
        shape.shape_type is not None and shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        for slide in prs.slides
        for shape in slide.shapes
    )
    assert has_picture, "expected at least one embedded chart image across all slides"
